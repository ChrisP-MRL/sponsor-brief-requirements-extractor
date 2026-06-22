#!/usr/bin/env python3
"""
Convert Microsoft PowerPoint (.pptx) files to plain text, organized by slides.
Requires: python-pptx library (install with: pip install python-pptx)
"""

import sys
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("Error: python-pptx library not found.")
    print("Install it with: pip install python-pptx")
    sys.exit(1)


def pptx_to_text(pptx_path, output_path=None):
    """
    Convert a PowerPoint presentation to plain text.

    Args:
        pptx_path: Path to the .pptx file
        output_path: Optional path for output text file. If None, prints to stdout.

    Returns:
        str: The extracted text content organized by slides
    """
    try:
        # Load the PowerPoint presentation
        prs = Presentation(pptx_path)

        # Extract text from all slides
        text_content = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            text_content.append(f"=== SLIDE {slide_num} ===\n")

            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text)
                    text_content.append("")  # Add blank line after each shape

            text_content.append("")  # Add blank line after each slide

        # Join all content
        full_text = '\n'.join(text_content)

        # Save or print the result
        if output_path:
            output_dir = Path(output_path).parent
            output_dir.mkdir(parents=True, exist_ok=True)

            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(full_text)
            print(f"Text extracted and saved to: {output_path}")
        else:
            print(full_text)

        return full_text

    except Exception as e:
        print(f"Error processing PowerPoint file: {e}")
        sys.exit(1)


def main():
    """Main entry point for the script."""
    if len(sys.argv) < 2:
        print("Usage: python pptx_to_text.py <input.pptx> [output.txt]")
        print("\nExamples:")
        print("  python pptx_to_text.py presentation.pptx           # Print to console")
        print("  python pptx_to_text.py presentation.pptx output.txt  # Save to file")
        sys.exit(1)

    input_file = Path(sys.argv[1])

    if not input_file.exists():
        print(f"Error: File not found: {input_file}")
        sys.exit(1)

    if not input_file.suffix.lower() == '.pptx':
        print("Warning: File does not have .pptx extension")

    output_file = sys.argv[2] if len(sys.argv) > 2 else None

    pptx_to_text(input_file, output_file)


if __name__ == "__main__":
    main()
